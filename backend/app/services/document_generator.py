import os
from typing import Tuple, Dict, Any, Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document as DocxDocument
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from app.core.config import settings
from app.services.model_manager import model_manager
from app.services.rag_service import rag_service
from app.services.tts_service import tts_service
import json
import logging
import zipfile
from io import BytesIO

logger = logging.getLogger(__name__)

class DocumentGenerator:
    """Service for generating documents from chat content"""
    
    async def generate(
        self,
        content: str,
        doc_type: str,
        user_context: Dict[str, Any],
        session_id: Optional[int] = None,
        topic: Optional[str] = None,
        template_id: Optional[int] = None
    ) -> Tuple[bytes, str, str]:
        """Generate a document and return file data, filename, and content type"""
        
        if doc_type == "ppt":
            return await self._generate_ppt(content, user_context, session_id, topic, template_id)
        elif doc_type == "mp4":
            # MP4 video generation should use HeyGen, not PowerPoint conversion
            return await self._generate_video_heygen(content, user_context, session_id, topic)
        elif doc_type == "doc":
            return await self._generate_docx(content, user_context)
        elif doc_type == "pdf":
            return await self._generate_pdf(content, user_context)
        elif doc_type == "podcast":
            # Explicit podcast/dialogue generation (two voices)
            return await self._generate_podcast(content, user_context, "mp3", session_id, topic)
        elif doc_type == "mp3" or doc_type == "wav" or doc_type == "speech":
            # Default to speech/monologue generation (single voice, no dialogue)
            # mp3/wav requests default to speech unless explicitly requesting podcast
            audio_format = "mp3" if doc_type == "speech" else doc_type
            return await self._generate_speech(content, user_context, session_id, topic, audio_format)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    async def _generate_ppt(self, content: str, user_context: Dict[str, Any], session_id: Optional[int] = None, topic: Optional[str] = None, template_id: Optional[int] = None) -> Tuple[bytes, str, str]:
        """Generate PowerPoint presentation using Presenton.ai API or fallback to local generation"""
        # Try Presenton.ai API first if configured
        if settings.PRESENTON_API_KEY:
            logger.info("=" * 60)
            logger.info("POWERPOINT GENERATION: Attempting Presenton.ai API...")
            logger.info(f"Content length: {len(content)} chars, Topic: {topic}, Template ID: {template_id}")
            logger.info("=" * 60)
            
            try:
                from app.services.presenton_service import presenton_service
                logger.info("✓ Presenton service imported successfully")
                
                # Try to find cisco_template.pptx
                template_path = None
                if template_id:
                    try:
                        from app.models.document import Document, DocumentType
                        from app.core.database import SessionLocal
                        db = SessionLocal()
                        try:
                            template_doc = db.query(Document).filter(
                                Document.id == template_id,
                                Document.file_type.in_([DocumentType.PPT, DocumentType.PPTX])
                            ).first()
                            if template_doc and template_doc.file_path:
                                template_path = template_doc.file_path
                                logger.info(f"✓ Found template: {template_path} (ID: {template_id})")
                        finally:
                            db.close()
                    except Exception as e:
                        logger.warning(f"Could not load template {template_id}: {e}", exc_info=True)
                
                # Use Presenton.ai API to generate PowerPoint
                logger.info("🚀 Calling Presenton.ai API...")
                presenton_result, filename = await presenton_service.generate_powerpoint(
                    content=content,
                    topic=topic,
                    template_path=template_path
                )
                
                logger.info("=" * 60)
                logger.info(f"✅ SUCCESS: Presenton.ai generated PowerPoint")
                logger.info(f"   Download path: {presenton_result.get('path')}")
                logger.info("=" * 60)
                
                # Return the path URL instead of bytes - frontend will download directly
                # Store the result in a way that frontend can access the path
                # We'll return a special format that the frontend can handle
                import json
                return json.dumps(presenton_result).encode('utf-8'), filename, "application/json"
                
            except ImportError:
                logger.warning("Presenton service not available, using local generation")
            except Exception as presenton_error:
                logger.warning(f"Presenton.ai API failed: {presenton_error}, falling back to local generation", exc_info=True)
        
        # Fallback: Local generation using python-pptx
        return await self._generate_ppt_local(content, user_context, session_id, topic, template_id)
    
    async def _generate_ppt_local(self, content: str, user_context: Dict[str, Any], session_id: Optional[int] = None, topic: Optional[str] = None, template_id: Optional[int] = None) -> Tuple[bytes, str, str]:
        """Generate PowerPoint presentation locally using python-pptx (fallback method)"""
        # Try to load custom template if template_id provided
        template_path = None
        if template_id:
            try:
                from app.models.document import Document, DocumentType
                from app.core.database import SessionLocal
                db = SessionLocal()
                try:
                    template_doc = db.query(Document).filter(
                        Document.id == template_id,
                        Document.file_type.in_([DocumentType.PPT, DocumentType.PPTX])
                    ).first()
                    if template_doc and template_doc.file_path:
                        template_path = template_doc.file_path
                        logger.info(f"Using custom template: {template_path} (ID: {template_id})")
                    else:
                        logger.warning(f"Template {template_id} not found or has no file path")
                finally:
                    db.close()
            except Exception as e:
                logger.warning(f"Could not load template {template_id}: {e}", exc_info=True)
        
        # Load template or create new presentation
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            logger.info("Loaded custom PowerPoint template")
            # Don't override dimensions if using template (preserve template settings)
        else:
            # Use default template or check for default template in uploads
            default_template_path = os.path.join(settings.UPLOAD_DIR, "templates", "default_template.pptx")
            if os.path.exists(default_template_path):
                prs = Presentation(default_template_path)
                logger.info("Using default template")
            else:
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                logger.info("Using blank presentation (no template)")
        
        # If using a template, try to fill existing slides; otherwise create new slides
        # For now, we'll add new slides to the template (preserving template design)
        # Check if template already has slides
        if len(prs.slides) > 0:
            # Template has slides - we can either:
            # 1. Fill existing placeholders in template slides
            # 2. Add new slides using template layouts
            # For now, we'll add new slides using template layouts
            logger.info(f"Template has {len(prs.slides)} existing slides. Adding new slides with content.")
        
        # Title slide (use first slide if template, or create new)
        if len(prs.slides) == 0:
            title_slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
            if title_slide_layout:
                slide = prs.slides.add_slide(title_slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = topic or "AI Generated Presentation"
                # Try to set subtitle if placeholder exists
                if len(slide.placeholders) > 1:
                    try:
                        slide.placeholders[1].text = f"Generated by {user_context.get('full_name', 'User')}"
                    except:
                        pass
        else:
            # Update first slide if it's a title slide
            first_slide = prs.slides[0]
            if first_slide.shapes.title:
                first_slide.shapes.title.text = topic or "AI Generated Presentation"
        
        # Split content into slides (by paragraphs or sections)
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        
        # If content is structured, use it; otherwise create slides from paragraphs
        # Limit to 10 slides as per requirement
        for i, para in enumerate(paragraphs[:10]):  # Limit to 10 slides
            if para:
                try:
                    # Try to get a bullet slide layout (index 1), fallback to first available layout
                    if len(prs.slide_layouts) > 1:
                        bullet_slide_layout = prs.slide_layouts[1]
                    elif len(prs.slide_layouts) > 0:
                        bullet_slide_layout = prs.slide_layouts[0]
                    else:
                        logger.warning("No slide layouts available, skipping slide creation")
                        continue
                    
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    
                    # Extract title from first line or use slide number
                    lines = para.split('\n')
                    title_text = lines[0][:60] if lines else f"Slide {i + 2}"
                    
                    # Try to set title if title shape exists
                    try:
                        if shapes.title:
                            shapes.title.text = title_text
                    except (AttributeError, IndexError) as title_error:
                        # No title shape available, try to find title placeholder
                        try:
                            for placeholder in shapes.placeholders:
                                if placeholder.placeholder_format.type == 0:  # Title placeholder type
                                    placeholder.text = title_text
                                    break
                        except:
                            logger.warning(f"Could not set title: {title_error}")
                    
                    # Try to find a text placeholder for body content
                    body_shape = None
                    for placeholder in shapes.placeholders:
                        if placeholder.placeholder_format.type == 1:  # Body placeholder type
                            body_shape = placeholder
                            break
                    
                    # If no body placeholder found, try to use placeholder at index 1, or any text placeholder
                    if not body_shape:
                        if len(shapes.placeholders) > 1:
                            try:
                                body_shape = shapes.placeholders[1]
                            except:
                                pass
                        if not body_shape:
                            # Try to find any text placeholder
                            for placeholder in shapes.placeholders:
                                try:
                                    if hasattr(placeholder, 'text_frame'):
                                        body_shape = placeholder
                                        break
                                except:
                                    continue
                    
                    # If we found a body shape, add content
                    if body_shape and hasattr(body_shape, 'text_frame'):
                        try:
                            tf = body_shape.text_frame
                            tf.text = ""
                            content_text = '\n'.join(lines[1:]) if len(lines) > 1 else para
                            # Split into bullet points if needed
                            bullets = [b.strip() for b in content_text.split('\n') if b.strip()][:6]
                            for j, bullet in enumerate(bullets):
                                if j == 0:
                                    if len(tf.paragraphs) > 0:
                                        p = tf.paragraphs[0]
                                        p.text = bullet[:200]
                                        p.level = 0
                                    else:
                                        p = tf.add_paragraph()
                                        p.text = bullet[:200]
                                        p.level = 0
                                else:
                                    p = tf.add_paragraph()
                                    p.text = bullet[:200]
                                    p.level = 0
                        except Exception as text_error:
                            logger.warning(f"Could not add text to slide: {text_error}")
                    else:
                        logger.warning(f"Could not find suitable placeholder for slide content")
                        
                except Exception as slide_error:
                    logger.error(f"Error creating slide {i + 2}: {slide_error}", exc_info=True)
                    # Continue with next slide instead of failing completely
                    continue
        
        # Save to bytes
        from io import BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.getvalue(), "presentation.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    
    async def _generate_docx(self, content: str, user_context: Dict[str, Any]) -> Tuple[bytes, str, str]:
        """Generate Word document"""
        doc = DocxDocument()
        
        # Add title
        doc.add_heading('AI Generated Document', 0)
        doc.add_paragraph(f'Generated by {user_context.get("full_name", "User")}')
        doc.add_paragraph('')
        
        # Add content
        paragraphs = content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                doc.add_paragraph(para)
                doc.add_paragraph('')
        
        # Save to bytes
        from io import BytesIO
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        
        return output.getvalue(), "document.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    
    async def _generate_pdf(self, content: str, user_context: Dict[str, Any]) -> Tuple[bytes, str, str]:
        """Generate PDF document"""
        from io import BytesIO
        output = BytesIO()
        
        doc = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title = Paragraph("AI Generated Document", styles['Title'])
        story.append(title)
        story.append(Spacer(1, 12))
        
        # Author
        author = Paragraph(f"Generated by {user_context.get('full_name', 'User')}", styles['Normal'])
        story.append(author)
        story.append(Spacer(1, 24))
        
        # Content
        paragraphs = content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                p = Paragraph(para[:1000], styles['Normal'])  # Limit length
                story.append(p)
                story.append(Spacer(1, 12))
        
        doc.build(story)
        output.seek(0)
        
        return output.getvalue(), "document.pdf", "application/pdf"
    
    async def _generate_video_heygen(self, content: str, user_context: Dict[str, Any], session_id: Optional[int] = None, topic: Optional[str] = None) -> Tuple[bytes, str, str]:
        """Generate MP4 video using HeyGen API"""
        logger.info("=" * 60)
        logger.info("VIDEO GENERATION: Using HeyGen API")
        logger.info(f"Content length: {len(content)} chars, Topic: {topic}")
        logger.info("=" * 60)
        
        try:
            from app.services.heygen_service import heygen_service
            
            # Format content as video script
            script = self._format_content_as_video_script(content, topic)
            
            # Generate video using HeyGen
            result = await heygen_service.generate_video(
                script=script,
                topic=topic
            )
            
            video_url = result.get("video_url")
            if not video_url:
                raise ValueError("HeyGen did not return a video URL")
            
            # Download the video from HeyGen URL
            import httpx
            async with httpx.AsyncClient(timeout=300.0) as client:
                logger.info(f"Downloading video from HeyGen: {video_url}")
                response = await client.get(video_url)
                response.raise_for_status()
                video_data = response.content
            
            filename = result.get("filename", f"video_{topic or 'generated'}.mp4")
            return video_data, filename, "video/mp4"
            
        except Exception as e:
            logger.error(f"HeyGen video generation failed: {e}", exc_info=True)
            raise ValueError(f"Failed to generate video using HeyGen: {str(e)}")
    
    def _format_content_as_video_script(self, content: str, topic: Optional[str] = None) -> str:
        """Format content as a video script suitable for narration"""
        import re
        
        # Remove markdown formatting, URLs, and citations
        script = content
        
        # Remove markdown links [text](url) -> text
        script = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', script)
        
        # Remove URLs
        script = re.sub(r'https?://[^\s]+', '', script)
        
        # Remove citation markers
        script = re.sub(r'\[(\d+)\]', '', script)
        script = re.sub(r'\(Source:[^\)]+\)', '', script)
        script = re.sub(r'Sources?:[^\n]+', '', script, flags=re.IGNORECASE)
        
        # Remove excessive line breaks
        script = re.sub(r'\n{3,}', '\n\n', script)
        
        # Clean up whitespace
        script = ' '.join(script.split())
        
        # Add introduction if topic provided
        if topic:
            intro = f"Today, we'll explore {topic}. "
            script = intro + script
        
        # Limit to approximately 3 minutes (450 words at 150 words/minute)
        words = script.split()
        if len(words) > 450:
            script = ' '.join(words[:450]) + "..."
        
        return script.strip()
    
    async def _convert_ppt_to_mp4(self, ppt_data: bytes, audio_data: Optional[bytes], topic: str) -> Tuple[bytes, str, str]:
        """Convert PowerPoint to MP4 video"""
        import tempfile
        import os
        
        try:
            # Try using LibreOffice headless (cross-platform)
            return await self._convert_ppt_to_mp4_libreoffice(ppt_data, audio_data, topic)
        except Exception as libre_error:
            logger.warning(f"LibreOffice conversion failed: {libre_error}")
        
        try:
            # Try using PowerPoint COM automation (Windows only)
            if os.name == 'nt':
                return await self._convert_ppt_to_mp4_powerpoint(ppt_data, audio_data, topic)
        except Exception as ppt_error:
            logger.warning(f"PowerPoint automation failed: {ppt_error}")
        
        # If all methods fail, raise to trigger package fallback
        raise ValueError("PPT to MP4 conversion not available. Install LibreOffice or use Windows with PowerPoint.")
    
    async def _convert_ppt_to_mp4_libreoffice(self, ppt_data: bytes, audio_data: Optional[bytes], topic: str) -> Tuple[bytes, str, str]:
        """Convert PPT to MP4 using LibreOffice headless"""
        import tempfile
        import subprocess
        import asyncio
        import os
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save PPT
            ppt_path = os.path.join(temp_dir, "presentation.pptx")
            with open(ppt_path, "wb") as f:
                f.write(ppt_data)
            
            # Try to convert using LibreOffice
            # Note: LibreOffice can export to PDF/images, but not directly to MP4
            # We'll need to use a different approach or combine with ffmpeg
            
            # For now, this is a placeholder - actual implementation would require
            # LibreOffice + ffmpeg or similar toolchain
            raise NotImplementedError("LibreOffice conversion requires additional setup")
    
    async def _convert_ppt_to_mp4_powerpoint(self, ppt_data: bytes, audio_data: Optional[bytes], topic: str) -> Tuple[bytes, str, str]:
        """Convert PPT to MP4 using PowerPoint COM automation (Windows only)"""
        import tempfile
        import win32com.client
        import pythoncom
        import os
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save PPT
            ppt_path = os.path.join(temp_dir, "presentation.pptx")
            with open(ppt_path, "wb") as f:
                f.write(ppt_data)
            
            # Initialize COM
            pythoncom.CoInitialize()
            try:
                # Open PowerPoint
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                
                # Open presentation
                presentation = powerpoint.Presentations.Open(ppt_path)
                
                # Export as MP4
                mp4_path = os.path.join(temp_dir, "output.mp4")
                presentation.ExportAsFixedFormat(
                    mp4_path,
                    FixedFormatType=26,  # ppFixedFormatTypeMP4
                    Intent=0,  # ppFixedFormatIntentScreen
                    FrameSlides=0,
                    HandoutOrder=0,
                    OutputType=0,
                    PrintHiddenSlides=0,
                    PrintRange=None,
                    RangeType=0,
                    SlideShowName="",
                    IncludeDocProperties=True,
                    KeepIRMSettings=True,
                    DocStructureTags=True,
                    BitmapMissingFonts=True,
                    UseISO19005_1=False
                )
                
                presentation.Close()
                powerpoint.Quit()
                
                # Read MP4 file
                with open(mp4_path, "rb") as f:
                    video_data = f.read()
                
                logger.info(f"Successfully converted PPT to MP4 using PowerPoint ({len(video_data)} bytes)")
                return video_data, "demo_video.mp4", "video/mp4"
                
            finally:
                pythoncom.CoUninitialize()
    
    async def _create_video_package(self, ppt_data: bytes, ppt_filename: str, audio_data: Optional[bytes], audio_filename: Optional[str], topic: Optional[str], user_context: Dict[str, Any]) -> Tuple[bytes, str, str]:
        """Create a ZIP package with PPT, Audio, and instructions"""
        from io import BytesIO
        import zipfile
        
        output = BytesIO()
        with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Add PowerPoint
            zip_file.writestr(ppt_filename, ppt_data)
            
            # Add Audio if available
            if audio_data and audio_filename:
                zip_file.writestr(audio_filename, audio_data)
            
            # Add instructions
            instructions = f"""VIDEO CREATION PACKAGE
Generated by: {user_context.get('full_name', 'User')}
Topic: {topic or 'Business Content'}

This package contains:
1. {ppt_filename} - PowerPoint presentation with slides
{f"2. {audio_filename} - Audio narration (MP3)" if audio_data else "2. Audio generation failed - please generate separately"}

TO CREATE YOUR VIDEO:

Option 1: Use PowerPoint (Recommended)
1. Open {ppt_filename} in PowerPoint
2. Go to File > Export > Create a Video
3. Select "Use Recorded Timings and Narrations"
4. If you have {audio_filename}, add it as background audio
5. Export as MP4

Option 2: Use Online Tools
- Upload {ppt_filename} to: https://www.online-convert.com/
- Or use: https://cloudconvert.com/pptx-to-mp4
- Add {audio_filename} as audio track if available

Option 3: Use AI Video Services
- Synthesia.io - AI video generation from text
- RunwayML - Text to video AI
- Lumen5 - Create videos from content

Option 4: Use Video Editing Software
- Import {ppt_filename} slides as images
- Add {audio_filename} as audio track
- Export as MP4

For professional results, consider using:
- Synthesia (AI avatars + voice)
- RunwayML (Text-to-video)
- Adobe Premiere Pro
- Camtasia
"""
            zip_file.writestr("VIDEO_CREATION_INSTRUCTIONS.txt", instructions.encode('utf-8'))
        
        output.seek(0)
        logger.info("Created video package with PPT and audio")
        return output.getvalue(), "video_package.zip", "application/zip"
    
    
    async def generate_podcast_script(
        self,
        content: str,
        user_context: Dict[str, Any],
        topic: Optional[str] = None
    ) -> str:
        """Generate only the podcast dialogue script (without audio)"""
        try:
            # Use LLM to generate a dialogue between 2 people
            dialogue_prompt = f"""Create a natural, engaging, and comprehensive podcast dialogue between two people by extracting ALL key information from the content below. 

CRITICAL INSTRUCTIONS:
1. **EXTRACT ALL KEY POINTS**: Dig deep and extract ALL important information, features, benefits, details, and examples from the content. Include specific details, not just summaries.

2. **BE COMPREHENSIVE**: Include all relevant information that would be valuable to listeners. Extract numbers, statistics, specific features, benefits, use cases, and technical details.

3. **SYNTHESIZE INFORMATION**: Combine related information from different parts of the content to create a cohesive discussion.

4. **INCLUDE ALL DETAILS**: Don't skip important details. If the content mentions specific features, benefits, or examples, include them all in the dialogue.

Make it conversational, informative, comprehensive, and easy to follow. Format it as a script with clear speaker labels.

Topic: {topic or 'Business Content'}
Content to extract information from (extract ALL key points):
{content[:2000]}

Format the dialogue like this:
[Host]: Welcome to today's podcast. Let's discuss...
[Guest]: Thanks for having me. I think...
[Host]: That's interesting. Can you elaborate on...
[Guest]: Absolutely. The key point is...

Continue the dialogue for about 10-15 exchanges, extracting and including ALL important information from the content. Make it informative, engaging, and comprehensive. End with a natural conclusion."""

            # Get the LLM model
            try:
                llm = model_manager.get_chat_model(model_id="auto", temperature=0.7)
            except Exception as e:
                logger.warning(f"Could not get LLM for dialogue generation: {e}")
                # Fallback: create a simple dialogue from content
                return self._create_simple_dialogue(content, topic)
            else:
                # Generate dialogue using LLM
                from langchain.schema import HumanMessage
                from langchain_openai import AzureChatOpenAI, ChatOpenAI
                from langchain_aws import ChatBedrock, BedrockLLM
                
                try:
                    if isinstance(llm, ChatBedrock):
                        # ChatBedrock with retry logic for throttling
                        import time
                        messages = [HumanMessage(content=dialogue_prompt)]
                        max_retries = 5
                        base_delay = 2
                        
                        for attempt in range(max_retries):
                            try:
                                # Run blocking LLM call in thread pool
                                response = await asyncio.to_thread(llm.invoke, messages)
                                break
                            except Exception as bedrock_error:
                                error_str = str(bedrock_error)
                                is_throttling = (
                                    "ThrottlingException" in error_str or 
                                    "Too many requests" in error_str or
                                    "throttl" in error_str.lower()
                                )
                                
                                if is_throttling and attempt < max_retries - 1:
                                    delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                    logger.warning(f"Bedrock throttling in dialogue generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                    await asyncio.sleep(delay)
                                    continue
                                
                                if is_throttling:
                                    raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                    raise
                    elif isinstance(llm, (AzureChatOpenAI, ChatOpenAI)):
                        messages = [HumanMessage(content=dialogue_prompt)]
                        # Add user parameter for Cisco if needed
                        invoke_kwargs = {}
                        if isinstance(llm, AzureChatOpenAI) and settings.CISCO_APPKEY:
                            user_data = {"appkey": settings.CISCO_APPKEY}
                            invoke_kwargs["user"] = json.dumps(user_data)
                        # Run blocking LLM call in thread pool
                        response = await asyncio.to_thread(llm.invoke, messages, **invoke_kwargs)
                    elif isinstance(llm, BedrockLLM):
                        # BedrockLLM with retry logic for throttling
                        import time
                        max_retries = 5
                        base_delay = 2
                        
                        for attempt in range(max_retries):
                            try:
                                # Run blocking LLM call in thread pool
                                response = await asyncio.to_thread(llm.invoke, dialogue_prompt)
                                break
                            except Exception as bedrock_error:
                                error_str = str(bedrock_error)
                                is_throttling = (
                                    "ThrottlingException" in error_str or 
                                    "Too many requests" in error_str or
                                    "throttl" in error_str.lower()
                                )
                                
                                if is_throttling and attempt < max_retries - 1:
                                    delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                    logger.warning(f"Bedrock throttling in dialogue generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                    await asyncio.sleep(delay)
                                    continue
                                
                                if is_throttling:
                                    raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                raise
                    elif isinstance(llm, ChatOpenAI):
                        # ChatOpenAI also uses messages format
                        messages = [HumanMessage(content=dialogue_prompt)]
                        # Run blocking LLM call in thread pool
                        response = await asyncio.to_thread(llm.invoke, messages)
                    else:
                        # Fallback for other model types
                        messages = [HumanMessage(content=dialogue_prompt)]
                        # Run blocking LLM call in thread pool
                        response = await asyncio.to_thread(llm.invoke, messages)
                    
                    # Extract dialogue text
                    if hasattr(response, 'content'):
                        dialogue = response.content
                    elif isinstance(response, str):
                        dialogue = response
                    else:
                        dialogue = str(response)
                    
                    logger.info(f"Generated podcast script ({len(dialogue)} characters)")
                    return dialogue
                except Exception as e:
                    logger.error(f"Error generating dialogue with LLM: {e}")
                    return self._create_simple_dialogue(content, topic)
        except Exception as e:
            logger.error(f"Error in podcast script generation: {e}")
            return self._create_simple_dialogue(content, topic)
    
    async def _generate_podcast(
        self, 
        content: str, 
        user_context: Dict[str, Any], 
        audio_format: str,
        session_id: Optional[int] = None,
        topic: Optional[str] = None
    ) -> Tuple[bytes, str, str]:
        """Generate a podcast dialogue script (and optionally audio) from content"""
        # Check if content is already a dialogue script
        is_script = "[Host]" in content or "[Guest]" in content or "[host]" in content or "[guest]" in content
        
        if is_script:
            # Content is already a script, use it directly
            dialogue = content
            logger.info("Content is already a podcast script, using it directly")
        else:
            # Generate dialogue script from content
            try:
                # Use LLM to generate a dialogue between 2 people
                dialogue_prompt = f"""Create a natural, engaging, and comprehensive podcast dialogue between two people by extracting ALL key information from the content below. 

CRITICAL INSTRUCTIONS:
1. **EXTRACT ALL KEY POINTS**: Dig deep and extract ALL important information, features, benefits, details, and examples from the content. Include specific details, not just summaries.

2. **BE COMPREHENSIVE**: Include all relevant information that would be valuable to listeners. Extract numbers, statistics, specific features, benefits, use cases, and technical details.

3. **SYNTHESIZE INFORMATION**: Combine related information from different parts of the content to create a cohesive discussion.

4. **INCLUDE ALL DETAILS**: Don't skip important details. If the content mentions specific features, benefits, or examples, include them all in the dialogue.

Make it conversational, informative, comprehensive, and easy to follow. Format it as a script with clear speaker labels.

Topic: {topic or 'Business Content'}
Content to extract information from (extract ALL key points):
{content[:2000]}

Format the dialogue like this:
[Host]: Welcome to today's podcast. Let's discuss...
[Guest]: Thanks for having me. I think...
[Host]: That's interesting. Can you elaborate on...
[Guest]: Absolutely. The key point is...

Continue the dialogue for about 10-15 exchanges, extracting and including ALL important information from the content. Make it informative, engaging, and comprehensive. End with a natural conclusion."""

                # Get the LLM model
                try:
                    llm = model_manager.get_chat_model(model_id="auto", temperature=0.7)
                except Exception as e:
                    logger.warning(f"Could not get LLM for dialogue generation: {e}")
                    # Fallback: create a simple dialogue from content
                    dialogue = self._create_simple_dialogue(content, topic)
                else:
                    # Generate dialogue using LLM
                    from langchain.schema import HumanMessage
                    from langchain_openai import AzureChatOpenAI, ChatOpenAI
                    from langchain_aws import ChatBedrock, BedrockLLM
                    
                    try:
                        # Run LLM calls in thread pool to avoid blocking event loop
                        import asyncio
                        
                        async def invoke_llm():
                            if isinstance(llm, ChatBedrock):
                                # ChatBedrock with retry logic for throttling
                                import time
                                messages = [HumanMessage(content=dialogue_prompt)]
                                max_retries = 5
                                base_delay = 2
                                
                                for attempt in range(max_retries):
                                    try:
                                        response = await asyncio.to_thread(llm.invoke, messages)
                                        return response
                                    except Exception as bedrock_error:
                                        error_str = str(bedrock_error)
                                        is_throttling = (
                                            "ThrottlingException" in error_str or 
                                            "Too many requests" in error_str or
                                            "throttl" in error_str.lower()
                                        )
                                        
                                        if is_throttling and attempt < max_retries - 1:
                                            delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                            logger.warning(f"Bedrock throttling in dialogue generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                            await asyncio.sleep(delay)
                                            continue
                                        
                                        if is_throttling:
                                            raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                        raise
                            elif isinstance(llm, (AzureChatOpenAI, ChatOpenAI)):
                                messages = [HumanMessage(content=dialogue_prompt)]
                                # Add user parameter for Cisco if needed
                                invoke_kwargs = {}
                                if isinstance(llm, AzureChatOpenAI) and settings.CISCO_APPKEY:
                                    user_data = {"appkey": settings.CISCO_APPKEY}
                                    invoke_kwargs["user"] = json.dumps(user_data)
                                return await asyncio.to_thread(llm.invoke, messages, **invoke_kwargs)
                            elif isinstance(llm, BedrockLLM):
                                # BedrockLLM with retry logic for throttling
                                import time
                                max_retries = 5
                                base_delay = 2
                                
                                for attempt in range(max_retries):
                                    try:
                                        response = await asyncio.to_thread(llm.invoke, dialogue_prompt)
                                        return response
                                    except Exception as bedrock_error:
                                        error_str = str(bedrock_error)
                                        is_throttling = (
                                            "ThrottlingException" in error_str or 
                                            "Too many requests" in error_str or
                                            "throttl" in error_str.lower()
                                        )
                                        
                                        if is_throttling and attempt < max_retries - 1:
                                            delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                            logger.warning(f"Bedrock throttling in dialogue generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                            await asyncio.sleep(delay)
                                            continue
                                        
                                        if is_throttling:
                                            raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                        raise
                            elif isinstance(llm, ChatOpenAI):
                                # ChatOpenAI also uses messages format
                                messages = [HumanMessage(content=dialogue_prompt)]
                                return await asyncio.to_thread(llm.invoke, messages)
                            else:
                                # Fallback for other model types
                                messages = [HumanMessage(content=dialogue_prompt)]
                                return await asyncio.to_thread(llm.invoke, messages)
                        
                        response = await invoke_llm()
                        
                        # Extract dialogue text
                        if hasattr(response, 'content'):
                            dialogue = response.content
                        elif isinstance(response, str):
                            dialogue = response
                        else:
                            dialogue = str(response)
                    except Exception as e:
                        logger.error(f"Error generating dialogue with LLM: {e}")
                        dialogue = self._create_simple_dialogue(content, topic)
            except Exception as e:
                logger.error(f"Error in podcast generation: {e}")
                dialogue = self._create_simple_dialogue(content, topic)
        
        # Try to generate actual audio using TTS service
        # Run TTS in a thread pool to avoid blocking the event loop
        try:
            from app.services.tts_service import tts_service
            import asyncio
            
            logger.info(f"Generating {audio_format.upper()} audio from dialogue using TTS service...")
            # Run blocking TTS call in thread pool to avoid freezing UI
            audio_data = await asyncio.to_thread(
                tts_service.text_to_speech,
                text=dialogue,
                audio_format=audio_format,
                language="en",
                slow=False
            )
            
            # Return audio file
            import time
            content_type = f"audio/{audio_format}"
            safe_topic = (topic or 'content').replace(' ', '_').replace('/', '_')[:50] if topic else 'content'
            filename = f"podcast_{safe_topic}_{int(time.time())}.{audio_format}"
            
            logger.info(f"Successfully generated {audio_format.upper()} audio file: {filename}")
            return audio_data, filename, content_type
            
        except Exception as tts_error:
            logger.warning(f"TTS generation failed: {tts_error}. Returning dialogue script as fallback.")
            # Fallback: return dialogue as text file
            from io import BytesIO
            import time
            output = BytesIO()
            
            dialogue_text = f"""PODCAST DIALOGUE SCRIPT
Generated by: {user_context.get('full_name', 'User')}
Topic: {topic or 'Business Content'}
Format: {audio_format.upper()}

{dialogue}

---
Note: Audio generation failed ({str(tts_error)}). This is the dialogue script.
To convert to audio manually, you can use:
- Google Text-to-Speech (gTTS)
- Amazon Polly
- Azure Cognitive Services
- OpenAI TTS API
"""
            
            output.write(dialogue_text.encode('utf-8'))
            output.seek(0)
            
            content_type = "text/plain"
            filename = f"podcast_dialogue_{audio_format}.txt"
            
            return output.getvalue(), filename, content_type
    
    def _create_simple_dialogue(self, content: str, topic: Optional[str] = None) -> str:
        """Create a simple dialogue from content when LLM is not available"""
        lines = content.split('\n')[:10]  # Take first 10 lines
        dialogue = f"[Host]: Welcome! Today we're discussing {topic or 'this important topic'}.\n\n"
        dialogue += "[Guest]: Thanks for having me. Let me share some key points.\n\n"
        
        for i, line in enumerate(lines):
            if line.strip():
                speaker = "[Host]" if i % 2 == 0 else "[Guest]"
                dialogue += f"{speaker}: {line.strip()[:200]}\n\n"
        
        dialogue += "[Host]: That's very insightful. Thank you for sharing.\n\n"
        dialogue += "[Guest]: My pleasure. I hope this was helpful.\n"
        
        return dialogue
    
    async def _generate_speech(
        self,
        content: str,
        user_context: Dict[str, Any],
        session_id: Optional[int] = None,
        topic: Optional[str] = None,
        audio_format: str = "mp3"
    ) -> Tuple[bytes, str, str]:
        """Generate a speech/monologue (single voice) from content"""
        import time
        
        try:
            # Use LLM to generate a well-structured speech/monologue
            speech_text = None
            
            try:
                llm = model_manager.get_chat_model(model_id="auto", temperature=0.7)
            except Exception as e:
                logger.warning(f"Could not get LLM for speech generation: {e}")
                # Fallback: use content directly as speech
                speech_text = content
            else:
                # Generate speech using LLM
                from langchain.schema import HumanMessage
                from langchain_openai import AzureChatOpenAI, ChatOpenAI
                from langchain_aws import ChatBedrock, BedrockLLM
                
                speech_prompt = f"""Create a compelling, comprehensive, and well-structured speech or monologue by extracting ALL key information from the content below. 

CRITICAL INSTRUCTIONS:
1. **EXTRACT ALL KEY POINTS**: Dig deep and extract ALL important information, features, benefits, details, and examples from the content. Don't just summarize - include specific details.

2. **BE COMPREHENSIVE**: Include all relevant information that would be valuable to the audience. Extract numbers, statistics, specific features, benefits, use cases, and technical details.

3. **SYNTHESIZE INFORMATION**: Combine related information from different parts of the content to create a cohesive narrative.

4. **INCLUDE ALL DETAILS**: Don't skip important details. If the content mentions specific features, benefits, or examples, include them all in the speech.

The speech should be:
- Clear, engaging, and comprehensive
- Well-organized with a natural flow
- Suitable for spoken delivery
- About 2-3 minutes when read aloud (approximately 300-500 words)
- Have a clear introduction, main points with all details, and conclusion

Topic: {topic or 'Business Content'}

Content to extract information from (extract ALL key points):
{content[:2000]}

Write the speech as a single, continuous monologue (no dialogue markers, no [Host] or [Guest] labels). Extract and include ALL important information from the content above. Make it sound natural and conversational, as if someone is giving a presentation or speech."""

                try:
                    import asyncio
                    
                    if isinstance(llm, ChatBedrock):
                        # ChatBedrock with retry logic for throttling
                        messages = [HumanMessage(content=speech_prompt)]
                        max_retries = 5
                        base_delay = 2
                        
                        for attempt in range(max_retries):
                            try:
                                # Run blocking LLM call in thread pool
                                response = await asyncio.to_thread(llm.invoke, messages)
                                break
                            except Exception as bedrock_error:
                                error_str = str(bedrock_error)
                                is_throttling = (
                                    "ThrottlingException" in error_str or 
                                    "Too many requests" in error_str or
                                    "throttl" in error_str.lower()
                                )
                                
                                if is_throttling and attempt < max_retries - 1:
                                    delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                    logger.warning(f"Bedrock throttling in speech generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                    await asyncio.sleep(delay)
                                    continue
                                
                                if is_throttling:
                                    raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                raise
                    elif isinstance(llm, BedrockLLM):
                        # BedrockLLM with retry logic for throttling
                        max_retries = 5
                        base_delay = 2
                        
                        for attempt in range(max_retries):
                            try:
                                # Run blocking LLM call in thread pool
                                response = await asyncio.to_thread(llm.invoke, speech_prompt)
                                break
                            except Exception as bedrock_error:
                                error_str = str(bedrock_error)
                                is_throttling = (
                                    "ThrottlingException" in error_str or 
                                    "Too many requests" in error_str or
                                    "throttl" in error_str.lower()
                                )
                                
                                if is_throttling and attempt < max_retries - 1:
                                    delay = base_delay * (2 ** attempt) + (attempt * 0.5)
                                    logger.warning(f"Bedrock throttling in speech generation (attempt {attempt + 1}/{max_retries}). Retrying in {delay:.1f}s...")
                                    await asyncio.sleep(delay)
                                    continue
                                
                                if is_throttling:
                                    raise ValueError(f"AWS Bedrock rate limit exceeded. Please wait 30-60 seconds and try again.")
                                raise
                    elif isinstance(llm, (AzureChatOpenAI, ChatOpenAI)):
                        messages = [HumanMessage(content=speech_prompt)]
                        # Add user parameter for Cisco if needed
                        invoke_kwargs = {}
                        if isinstance(llm, AzureChatOpenAI) and settings.CISCO_APPKEY:
                            user_data = {"appkey": settings.CISCO_APPKEY}
                            invoke_kwargs["user"] = json.dumps(user_data)
                        # Run blocking LLM call in thread pool
                        response = await asyncio.to_thread(llm.invoke, messages, **invoke_kwargs)
                    else:
                        messages = [HumanMessage(content=speech_prompt)]
                        # Run blocking LLM call in thread pool
                        response = await asyncio.to_thread(llm.invoke, messages)
                    
                    # Extract speech text
                    if hasattr(response, 'content'):
                        speech_text = response.content
                    elif isinstance(response, str):
                        speech_text = response
                    else:
                        speech_text = str(response)
                    
                    logger.info(f"Generated speech text ({len(speech_text)} characters)")
                except Exception as e:
                    logger.error(f"Error generating speech with LLM: {e}", exc_info=True)
                    # Fallback: use content directly
                    speech_text = content
            
            # If speech_text is still None, use content as fallback
            if not speech_text:
                speech_text = content
            
            # Generate audio using TTS service (single voice, no dialogue parsing)
            # Run TTS in a thread pool to avoid blocking the event loop
            from app.services.tts_service import tts_service
            import asyncio
            
            logger.info(f"Generating {audio_format.upper()} audio from speech using TTS service...")
            # Run blocking TTS call in thread pool to avoid freezing UI
            audio_data = await asyncio.to_thread(
                tts_service.text_to_speech_dialogue,
                speech_text,
                audio_format=audio_format,
                language="en",
                use_dialogue=False  # Single voice, no dialogue parsing
            )
            
            # Return audio file
            content_type = f"audio/{audio_format}"
            safe_topic = (topic or 'speech').replace(' ', '_').replace('/', '_')[:50] if topic else 'speech'
            filename = f"speech_{safe_topic}_{int(time.time())}.{audio_format}"
            
            logger.info(f"Successfully generated {audio_format.upper()} speech audio file: {filename}")
            return audio_data, filename, content_type
            
        except Exception as tts_error:
            logger.warning(f"TTS generation failed: {tts_error}. Returning speech script as fallback.")
            # Fallback: return speech as text file
            from io import BytesIO
            import time
            output = BytesIO()
            
            speech_text_output = f"""SPEECH SCRIPT
Generated by: {user_context.get('full_name', 'User')}
Topic: {topic or 'Business Content'}
Format: {audio_format.upper()} Speech

{speech_text if 'speech_text' in locals() else content}

---
Note: Audio generation failed ({str(tts_error)}). This is the speech script.
To convert to audio manually, you can use:
- OpenAI TTS API
- Google Text-to-Speech (gTTS)
- Amazon Polly
- Azure Cognitive Services
"""
            
            output.write(speech_text_output.encode('utf-8'))
            output.seek(0)
            
            content_type = "text/plain"
            filename = f"speech_script_{int(time.time())}.txt"
            
            return output.getvalue(), filename, content_type

