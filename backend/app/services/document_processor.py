import os
from typing import Dict, Optional
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pydub import AudioSegment
from moviepy.editor import VideoFileClip
from app.core.config import settings

# Optional imports for audio/video transcription
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False
    whisper = None

class DocumentProcessor:
    def __init__(self):
        self.whisper_model = None
    
    def _load_whisper_model(self):
        """Lazy load Whisper model for audio/video transcription"""
        if not WHISPER_AVAILABLE:
            raise ImportError("Whisper is not available. Install with: pip install openai-whisper")
        if self.whisper_model is None:
            self.whisper_model = whisper.load_model("base")
        return self.whisper_model
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from various file types"""
        file_type_lower = file_type.lower()
        
        if file_type_lower == "pdf":
            return self._extract_from_pdf(file_path)
        elif file_type_lower in ["doc", "docx"]:
            return self._extract_from_docx(file_path)
        elif file_type_lower in ["xls", "xlsx"]:
            return self._extract_from_xlsx(file_path)
        elif file_type_lower in ["ppt", "pptx"]:
            return self._extract_from_pptx(file_path)
        elif file_type_lower in ["mp3", "wav", "m4a"]:
            return self._extract_from_audio(file_path)
        elif file_type_lower in ["mp4", "mov", "avi"]:
            return self._extract_from_video(file_path)
        elif file_type_lower in ["txt", "md"]:
            return self._extract_from_text(file_path)
        elif file_type_lower == "json":
            return self._extract_from_json(file_path)
        elif file_type_lower == "jsonl":
            return self._extract_from_jsonl(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = DocxDocument(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel files"""
        workbook = load_workbook(file_path)
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}\n")
            for row in sheet.iter_rows(values_only=True):
                text_parts.append("\t".join([str(cell) if cell else "" for cell in row]))
        return "\n".join(text_parts)
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {slide_num}:\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    
    def _extract_from_audio(self, file_path: str) -> str:
        """Extract text from audio using Whisper"""
        model = self._load_whisper_model()
        result = model.transcribe(file_path)
        return result["text"]
    
    def _extract_from_video(self, file_path: str) -> str:
        """Extract text from video (audio track) using Whisper"""
        # Extract audio from video
        video = VideoFileClip(file_path)
        audio_path = file_path.replace(Path(file_path).suffix, ".wav")
        video.audio.write_audiofile(audio_path, verbose=False, logger=None)
        video.close()
        
        # Transcribe audio
        text = self._extract_from_audio(audio_path)
        
        # Clean up temporary audio file
        if os.path.exists(audio_path):
            os.remove(audio_path)
        
        return text
    
    def _extract_from_text(self, file_path: str) -> str:
        """Extract text from plain text files"""
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    
    def _extract_from_json(self, file_path: str) -> str:
        """Extract text from JSON files"""
        import json
        
        # Read file in binary mode first to check for BOM
        with open(file_path, "rb") as f:
            raw_data = f.read()
        
        # Detect BOM and set appropriate encoding
        if raw_data.startswith(b'\xff\xfe'):
            # UTF-16 LE BOM
            encoding = 'utf-16-le'
            raw_data = raw_data[2:]  # Remove BOM
        elif raw_data.startswith(b'\xfe\xff'):
            # UTF-16 BE BOM
            encoding = 'utf-16-be'
            raw_data = raw_data[2:]  # Remove BOM
        elif raw_data.startswith(b'\xef\xbb\xbf'):
            # UTF-8 BOM
            encoding = 'utf-8-sig'
            raw_data = raw_data[3:]  # Remove BOM
        else:
            encoding = None
        
        # Try multiple encodings - UTF-16 is common for Windows-generated JSON files
        encodings_to_try = []
        if encoding:
            encodings_to_try.append(encoding)
        encodings_to_try.extend(['utf-8', 'utf-8-sig', 'utf-16', 'utf-16-le', 'utf-16-be', 'latin-1', 'cp1252'])
        
        # Try each encoding until one works
        last_error = None
        for enc in encodings_to_try:
            try:
                # If we detected BOM and this is the detected encoding, use the data we already read
                if enc == encoding and encoding:
                    # Use the data we already read (without BOM)
                    text = raw_data.decode(enc, errors='strict')
                else:
                    # Read file fresh with this encoding
                    with open(file_path, "r", encoding=enc) as file:
                        text = file.read()
                
                # Try to parse JSON
                data = json.loads(text)
                
                # Convert JSON to readable text format
                # For simple JSON objects, convert to key-value pairs
                if isinstance(data, dict):
                    text_parts = []
                    for key, value in data.items():
                        if isinstance(value, (dict, list)):
                            text_parts.append(f"{key}: {json.dumps(value, indent=2)}")
                        else:
                            text_parts.append(f"{key}: {value}")
                    return "\n".join(text_parts)
                elif isinstance(data, list):
                    # For JSON arrays, format each item
                    return "\n".join([json.dumps(item, indent=2) for item in data])
                else:
                    return str(data)
            except (UnicodeDecodeError, UnicodeError) as e:
                last_error = e
                continue
            except json.JSONDecodeError as e:
                # If JSON decode fails, try next encoding
                last_error = e
                continue
        
        # Last resort: try to decode with errors='ignore' (better than 'replace' for JSON)
        for fallback_enc in ['utf-8', 'utf-16-le', 'utf-16-be', 'latin-1']:
            try:
                text = raw_data.decode(fallback_enc, errors='ignore')
                if text.strip() and len(text) > 10:  # Only try if we got substantial text
                    data = json.loads(text)
                    if isinstance(data, dict):
                        text_parts = []
                        for key, value in data.items():
                            if isinstance(value, (dict, list)):
                                text_parts.append(f"{key}: {json.dumps(value, indent=2)}")
                            else:
                                text_parts.append(f"{key}: {value}")
                        return "\n".join(text_parts)
                    elif isinstance(data, list):
                        return "\n".join([json.dumps(item, indent=2) for item in data])
                    else:
                        return str(data)
            except (json.JSONDecodeError, ValueError):
                continue
            except Exception:
                continue
        
        # If all else fails, raise a helpful error
        raise ValueError(f"Could not decode JSON file with any encoding. The file may be corrupted or in an unsupported encoding. Last error: {last_error}")
    
    def _extract_from_jsonl(self, file_path: str) -> str:
        """Extract text from JSONL files (JSON Lines)"""
        import json
        
        # Try multiple encodings - UTF-16 is common for Windows-generated JSON files
        encodings_to_try = ['utf-8', 'utf-8-sig', 'utf-16', 'utf-16-le', 'utf-16-be', 'latin-1', 'cp1252']
        
        # Try each encoding until one works
        last_error = None
        for encoding in encodings_to_try:
            try:
                text_parts = []
                with open(file_path, "r", encoding=encoding) as file:
                    for line_num, line in enumerate(file, 1):
                        line = line.strip()
                        if not line:
                            continue
                        try:
                            data = json.loads(line)
                            # Format each JSON object
                            if isinstance(data, dict):
                                text_parts.append(f"Line {line_num}: {json.dumps(data, indent=2)}")
                            else:
                                text_parts.append(f"Line {line_num}: {str(data)}")
                        except json.JSONDecodeError:
                            # If line is not valid JSON, include it as plain text
                            text_parts.append(f"Line {line_num}: {line}")
                return "\n".join(text_parts)
            except (UnicodeDecodeError, UnicodeError) as e:
                last_error = e
                continue
        
        # If all encodings fail, try reading as binary and decode with errors='replace'
        try:
            text_parts = []
            with open(file_path, "rb") as file:
                for line_num, line_bytes in enumerate(file, 1):
                    line = line_bytes.decode('utf-8', errors='replace').strip()
                    if not line:
                        continue
                    try:
                        data = json.loads(line)
                        if isinstance(data, dict):
                            text_parts.append(f"Line {line_num}: {json.dumps(data, indent=2)}")
                        else:
                            text_parts.append(f"Line {line_num}: {str(data)}")
                    except json.JSONDecodeError:
                        text_parts.append(f"Line {line_num}: {line}")
            return "\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Could not decode JSONL file with any encoding. Last error: {last_error or e}")

# Global instance
document_processor = DocumentProcessor()

